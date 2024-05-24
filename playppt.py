from pptx import Presentation
from pptx.util import Inches

# 創建一個新的PPT
prs = Presentation()

# 添加標題幻燈片
slide_layout = prs.slide_layouts[0]  # 0是標題幻燈片的布局
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello, World!"
subtitle.text = "python-pptx was here!"

# 添加內容幻燈片
slide_layout = prs.slide_layouts[1]  # 1是標題和內容的布局
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Second Slide"
content.text = "This is the second slide content."

# 添加圖片幻燈片
slide_layout = prs.slide_layouts[5]  # 5是標題加圖片的布局
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Slide with Image"

img_path = 'static/futune_telling_logo.png'
left = Inches(1)
top = Inches(1.5)
height = Inches(3.5)
pic = slide.shapes.add_picture(img_path, left, top, height=height)

# 保存PPT
prs.save('test_presentation.pptx')

print("PPT created successfully.")
